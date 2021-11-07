import os

from pptx import Presentation


TEXT_NEED_REPLACE = [('圆', '锐')]


def replace_text(text_frame):  # 该函数实现的是文本替换功能
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            for tt in TEXT_NEED_REPLACE:
                if tt[0] in run.text:
                    run.text = run.text.replace(tt[0], tt[1])


def process_ppt(filename_open, filename_save):
    prs = Presentation(filename_open)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:  # 判断Shape是否含有文本框
                text_frame = shape.text_frame
                replace_text(text_frame)  # 调用replace_text函数实现文本替换
            if shape.has_table:  # 判断Shape是否含有表格
                table = shape.table
                for cell in table.iter_cells():  # 遍历表格的cell
                    text_frame = cell.text_frame
                    replace_text(text_frame)  # 调用replace_text函数实现文本替换
    prs.save(filename_save)  # 保存


def running(dirpath=os.path.dirname(__file__)):
    t = os.listdir(dirpath)
    a = []
    for i in range(len(t)):
        if t[i][-5:] == ".pptx" or t[i][-4:] == ".ppt" or t[i][-5:] == ".PPTX" or t[i][-4:] == ".PPT":
            a.append(t[i])
    # print(a)
    for i in a:
        if os.path.isfile("%s/%s" % (dirpath, i)):
            process_ppt("%s/%s" % (dirpath, i), "%s/%s" % (dirpath, i))
        elif os.path.isdir("%s/%s" % (dirpath, i)):
            running("%s/%s" % (dirpath, i))


if __name__ == "__main__":
    # print(__file__)
    running()